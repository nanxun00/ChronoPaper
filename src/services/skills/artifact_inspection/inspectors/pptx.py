"""PPTX 产物解析与质检。"""
from __future__ import annotations

from pathlib import Path

from src.services.skills.artifact_inspection.models import ArtifactInspectionResult
from src.services.skills.pptx_quality import (
    min_required_slides,
    parse_target_slide_count,
)

_PLACEHOLDER_PHRASES = (
    "方法1：具体",
    "方法2：创新",
    "本文研究的主要问题和意义",
    "当前领域存在的挑战或空白",
    "本文采用的主要技术路线",
    "主要实验发现",
    "研究的主要目标",
)


def _slide_texts(prs) -> list[tuple[str, list[str]]]:
    slides: list[tuple[str, list[str]]] = []
    for slide in prs.slides:
        title = ""
        bullets: list[str] = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text = shape.text_frame.text.strip()
            if not text:
                continue
            if not title:
                title = text.split("\n", 1)[0][:120]
            for line in text.split("\n"):
                line = line.strip()
                if line and line != title:
                    bullets.append(line[:200])
        slides.append((title or "（无标题）", bullets[:6]))
    return slides


def inspect_pptx(path: Path, *, query: str = "") -> ArtifactInspectionResult:
    rel = path.as_posix()
    errors: list[str] = []
    warnings: list[str] = []
    summary: list[str] = []

    if not path.is_file():
        return ArtifactInspectionResult(
            path=rel,
            kind="pptx",
            ok=False,
            errors=[f"文件不存在：{rel}"],
        )

    try:
        from pptx import Presentation

        prs = Presentation(str(path))
    except Exception as exc:
        return ArtifactInspectionResult(
            path=rel,
            kind="pptx",
            ok=False,
            errors=[f"无法读取 PPTX：{exc}"],
        )

    count = len(prs.slides)
    target = parse_target_slide_count(query)
    minimum = min_required_slides(target)
    summary.append(f"共 {count} 张幻灯片（目标约 {target} 张，至少 {minimum} 张）")

    if count < minimum:
        errors.append(
            f"仅 {count} 张，少于 {minimum} 张；需补全背景/方法/结果/讨论/结论等结构化页面"
        )

    slide_data = _slide_texts(prs)
    for idx, (title, bullets) in enumerate(slide_data[:8], start=1):
        preview = "；".join(bullets[:3]) if bullets else "（无要点）"
        summary.append(f"第{idx}页「{title}」：{preview[:160]}")

    all_text = "\n".join(title + "\n" + "\n".join(bullets) for title, bullets in slide_data)
    placeholder_hits = [p for p in _PLACEHOLDER_PHRASES if p in all_text]
    if len(placeholder_hits) >= 2:
        errors.append(
            "检测到模板占位文案（"
            + "、".join(placeholder_hits[:3])
            + "）；须改为论文具体内容"
        )
    elif placeholder_hits:
        warnings.append(f"疑似占位句：{placeholder_hits[0]}")

    if count >= minimum and sum(1 for _, b in slide_data if len(b) >= 2) < max(3, count // 3):
        warnings.append("多数页面 bullet 过少，内容可能过于简陋")

    return ArtifactInspectionResult(
        path=rel,
        kind="pptx",
        ok=not errors,
        errors=errors,
        warnings=warnings,
        summary_lines=summary,
    )
