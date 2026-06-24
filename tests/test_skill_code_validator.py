"""生成脚本校验单元测试。"""
from __future__ import annotations

import unittest

from src.services.skills.code_validator import validate_generated_code


class SkillCodeValidatorTests(unittest.TestCase):
    def test_allow_simple_pptx_script(self) -> None:
        code = """
from pathlib import Path
from pptx import Presentation

out = Path("output")
out.mkdir(parents=True, exist_ok=True)
prs = Presentation()
prs.slides.add_slide(prs.slide_layouts[0])
prs.save(out / "demo.pptx")
"""
        result = validate_generated_code(code)
        self.assertTrue(result.ok, result.errors)

    def test_reject_subprocess(self) -> None:
        code = "import subprocess\nsubprocess.run(['ls'])"
        result = validate_generated_code(code)
        self.assertFalse(result.ok)
        self.assertTrue(any("subprocess" in e for e in result.errors))

    def test_reject_os_system(self) -> None:
        code = "import os\nos.system('echo hi')"
        result = validate_generated_code(code)
        self.assertFalse(result.ok)

    def test_reject_eval(self) -> None:
        code = 'eval("1+1")'
        result = validate_generated_code(code)
        self.assertFalse(result.ok)

    def test_allow_tempfile_import(self) -> None:
        code = "import tempfile\nfrom pathlib import Path\nPath('output').mkdir(exist_ok=True)"
        result = validate_generated_code(code)
        self.assertTrue(result.ok, result.errors)

    def test_reviewable_shutil_rmtree(self) -> None:
        code = """
import shutil
from pathlib import Path
out = Path("output/runs/run1")
out.mkdir(parents=True, exist_ok=True)
shutil.rmtree(out / "tmp", ignore_errors=True)
"""
        result = validate_generated_code(code)
        self.assertFalse(result.ok, result.errors)
        self.assertTrue(result.needs_user_approval, result.errors)
        self.assertTrue(any("shutil" in e for e in result.reviewable_errors or []))

    def test_allow_reviewable_when_user_approved(self) -> None:
        code = "import shutil\nshutil.rmtree('output/runs/x/tmp', ignore_errors=True)"
        result = validate_generated_code(code, allow_reviewable=True)
        self.assertTrue(result.ok, result.errors)

    def test_reject_syntax_error(self) -> None:
        result = validate_generated_code("def foo(:\n  pass")
        self.assertFalse(result.ok)
        self.assertTrue(any("语法" in e for e in result.errors))

    def test_reject_fullpage_pdf_screenshot_deck(self) -> None:
        code = """
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches

pngs = [Path(f"output/assets/figures/doi_x_p{p:03d}.png") for p in range(1, 10)]
prs = Presentation()
for png in pngs:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_picture(str(png), Inches(0), Inches(0), width=Inches(13.333), height=Inches(7.5))
prs.save("output/runs/run1/final_presentation_cn.pptx")
"""
        result = validate_generated_code(code)
        self.assertFalse(result.ok)
        self.assertTrue(any("整页" in e for e in result.errors))


if __name__ == "__main__":
    unittest.main()
