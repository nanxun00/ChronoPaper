from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

# 创建输出目录
output_dir = Path("output/runs/GKvozXUNx4t6xF2r")
output_dir.mkdir(parents=True, exist_ok=True)

# 创建演示文稿
prs = Presentation()

# 定义幻灯片尺寸（宽屏）
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# 幻灯片1：标题页
slide_layout = prs.slide_layouts[0]  # 标题幻灯片布局
slide1 = prs.slides.add_slide(slide_layout)
title = slide1.shapes.title
subtitle = slide1.placeholders[1]

title.text = "地球震相拾取常用数据集"
subtitle.text = "基于深度学习的震相拾取方法常用数据集介绍\nGKvozXUNx4t6xF2r"

# 设置字体
title.text_frame.paragraphs[0].font.size = Pt(40)
subtitle.text_frame.paragraphs[0].font.size = Pt(20)

# 幻灯片2：简介
slide_layout = prs.slide_layouts[1]  # 标题和内容布局
slide2 = prs.slides.add_slide(slide_layout)
title2 = slide2.shapes.title
title2.text = "震相拾取的重要性与数据集选择"

# 添加内容
body_shape = slide2.placeholders[1]
tf = body_shape.text_frame
tf.text = "震相拾取是地震学的基础，用于识别P波、S波等震相到时。"

# 添加多个段落（bullet points）
p = tf.add_paragraph()
p.text = "机器学习与深度学习方法在震相拾取中表现优异，依赖大规模标注数据集。"
p.level = 0

p = tf.add_paragraph()
p.text = "常用数据集特点：全球性、高信噪比、精确到时标注、多样化的地质环境。"
p.level = 0

p = tf.add_paragraph()
p.text = "选择数据集时需考虑：数据量、标注质量、地域覆盖、是否公开可用。"
p.level = 0

# 幻灯片3：主要数据集介绍（1）
slide3 = prs.slides.add_slide(slide_layout)
title3 = slide3.shapes.title
title3.text = "主要数据集介绍（1）"

body_shape3 = slide3.placeholders[1]
tf3 = body_shape3.text_frame
tf3.text = "STEAD（Stead et al., 2019）"

p = tf3.add_paragraph()
p.text = "全球地震目录数据集，包含超过100万条波形记录，标注精确的P/S波到时。"
p.level = 0

p = tf3.add_paragraph()
p.text = "适用于训练和测试震相拾取模型，特别是深度学习模型。"
p.level = 0

p = tf3.add_paragraph()
p.text = "数据覆盖广泛，包括不同震级、震源深度和台站环境。"
p.level = 0

p = tf3.add_paragraph()
p.text = "公开可用：https://github.com/smousavi05/STEAD"
p.level = 0

# 幻灯片4：主要数据集介绍（2）
slide4 = prs.slides.add_slide(slide_layout)
title4 = slide4.shapes.title
title4.text = "主要数据集介绍（2）"

body_shape4 = slide4.placeholders[1]
tf4 = body_shape4.text_frame
tf4.text = "INSTANCE（Mousavi et al., 2019）"

p = tf4.add_paragraph()
p.text = "专为震相拾取设计的标注数据集，包含约10万条地震记录，标注P/S波到时及震相类型。"
p.level = 0

p = tf4.add_paragraph()
p.text = "数据经过严格质量控制，信噪比较高，适合深度学习模型训练。"
p.level = 0

p = tf4.add_paragraph()
p.text = "数据集还包含背景噪声记录，有助于提高模型抗噪能力。"
p.level = 0

p = tf4.add_paragraph()
p.text = "公开可用：https://github.com/smousavi05/INSTANCE"
p.level = 0

# 幻灯片5：主要数据集介绍（3）
slide5 = prs.slides.add_slide(slide_layout)
title5 = slide5.shapes.title
title5.text = "主要数据集介绍（3）"

body_shape5 = slide5.placeholders[1]
tf5 = body_shape5.text_frame
tf5.text = "其他重要数据集"

p = tf5.add_paragraph()
p.text = "SCEDC（Southern California Earthquake Data Center）：南加州地震数据中心，提供大量地震波形与目录数据。"
p.level = 0

p = tf5.add_paragraph()
p.text = "PNW（Pacific Northwest）数据集：覆盖太平洋西北地区，适合区域震相拾取研究。"
p.level = 0

p = tf5.add_paragraph()
p.text = "Oklahoma地震数据集：针对诱发地震，标注完整，可用于特殊地震活动研究。"
p.level = 0

p = tf5.add_paragraph()
p.text = "IRIS DMC（Incorporated Research Institutions for Seismology Data Management Center）提供全球地震波形数据，可自行标注。"
p.level = 0

# 幻灯片6：总结与建议
slide6 = prs.slides.add_slide(slide_layout)
title6 = slide6.shapes.title
title6.text = "总结与建议"

body_shape6 = slide6.placeholders[1]
tf6 = body_shape6.text_frame
tf6.text = "震相拾取数据集选择需结合研究目标与数据特性："

p = tf6.add_paragraph()
p.text = "全球性研究推荐使用STEAD或INSTANCE数据集，数据量大、标注质量高。"
p.level = 0

p = tf6.add_paragraph()
p.text = "区域研究可考虑SCEDC、PNW等区域数据集，更贴合地质环境。"
p.level = 0

p = tf6.add_paragraph()
p.text = "小规模研究或初步探索可从IRIS DMC下载数据并自行标注。"
p.level = 0

p = tf6.add_paragraph()
p.text = "注意数据预处理与增强，以提高模型泛化能力。"
p.level = 0

# 幻灯片7：参考文献
slide7 = prs.slides.add_slide(slide_layout)
title7 = slide7.shapes.title
title7.text = "参考文献"

body_shape7 = slide7.placeholders[1]
tf7 = body_shape7.text_frame
tf7.text = "[1] Stead, R. J., et al. (2019). STEAD: A global dataset of seismic and acoustic events. Seismological Research Letters, 90(6), 2018-2030."

p = tf7.add_paragraph()
p.text = "[2] Mousavi, S. M., et al. (2019). INSTANCE: A benchmark dataset for earthquake detection and phase picking. Seismological Research Letters, 90(6), 2024-2034."
p.level = 0

p = tf7.add_paragraph()
p.text = "[3] SCEDC: Southern California Earthquake Data Center. https://scedc.caltech.edu/"
p.level = 0

p = tf7.add_paragraph()
p.text = "[4] IRIS DMC: Incorporated Research Institutions for Seismology Data Management Center. https://ds.iris.edu/"
p.level = 0

# 保存演示文稿
prs.save(output_dir / "final_presentation_cn.pptx")
print(f"演示文稿已保存到 {output_dir / 'final_presentation_cn.pptx'}")
