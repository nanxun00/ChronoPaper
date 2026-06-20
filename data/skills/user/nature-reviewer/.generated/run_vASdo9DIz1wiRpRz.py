import pathlib
import fitz  # PyMuPDF
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import re
import json
import os
import sys

# 定义路径
ROOT_DIR = pathlib.Path(__file__).parent
PDF_PATH = ROOT_DIR / "input/papers/doi_10.1016_j.neucom.2023.126295/paper.pdf"
OUTPUT_DIR = ROOT_DIR / "output/runs/vASdo9DIz1wiRpRz"
OUTPUT_DIR.mkdir(parents=True, exist_ok=True)

# 输出PPTX路径
PPTX_PATH = OUTPUT_DIR / "innovation_review.pptx"

# 提取论文文本（使用PyMuPDF读取全文）
def extract_paper_text(pdf_path):
    doc = fitz.open(pdf_path)
    full_text = []
    for page in doc:
        full_text.append(page.get_text())
    return "\n".join(full_text)

# 分析创新点
def analyze_innovation(text, innovation_point):
    """分析文本中关于创新点的内容，返回相关段落和关键句"""
    # 这里我们使用简单的关键词匹配来定位创新点相关内容
    # 在实际应用中，可能需要更复杂的NLP处理
    
    patterns = {
        "cascade_dual_task": [r"cascade", r"dual.task", r"multi.task", r"级联双任务"],
        "perceptual_loss": [r"perceptual.loss", r"感知损失"],
        "error_consistency": [r"error.prediction.consistency", r"误差预测一致性"]
    }
    
    key_sentences = []
    paragraphs = text.split("\n\n")
    
    for paragraph in paragraphs:
        if any(re.search(pattern, paragraph.lower()) for pattern in patterns[innovation_point]):
            # 提取包含关键词的句子
            sentences = re.split(r'(?<=[.!?])\s+', paragraph)
            for sentence in sentences:
                if any(re.search(pattern, sentence.lower()) for pattern in patterns[innovation_point]):
                    # 清理句子
                    clean_sentence = re.sub(r'\s+', ' ', sentence).strip()
                    if len(clean_sentence) > 20:  # 过滤太短的句子
                        key_sentences.append(clean_sentence)
    
    return key_sentences

# 创建PPT
def create_ppt(paper_text):
    prs = Presentation()
    
    # 设置幻灯片尺寸为宽屏
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # 定义颜色
    title_color = RGBColor(0, 51, 102)  # 深蓝色
    text_color = RGBColor(0, 0, 0)
    highlight_color = RGBColor(178, 34, 34)  # 红色高亮
    
    # 第一页：标题页
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide1.shapes.title
    subtitle = slide1.placeholders[1]
    
    title.text = "论文创新点评审报告"
    title.text_frame.paragraphs[0].font.size = Pt(36)
    title.text_frame.paragraphs[0].font.color.rgb = title_color
    title.text_frame.paragraphs[0].font.bold = True
    
    subtitle.text = "论文：级联双任务、感知损失、误差预测一致性创新点分析"
    subtitle.text_frame.paragraphs[0].font.size = Pt(24)
    subtitle.text_frame.paragraphs[0].font.color.rgb = text_color
    
    # 添加运行ID到备注
    notes_slide1 = slide1.notes_slide
    notes_slide1.notes_text_frame.text = "运行ID: vASdo9DIz1wiRpRz"
    
    # 分析各创新点
    innovations = {
        "cascade_dual_task": "级联双任务创新",
        "perceptual_loss": "感知损失创新",
        "error_consistency": "误差预测一致性创新"
    }
    
    for innovation_key, innovation_title in innovations.items():
        # 为每个创新点创建幻灯片
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        
        # 设置标题
        slide.shapes.title.text = innovation_title
        slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(28)
        slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = title_color
        slide.shapes.title.text_frame.paragraphs[0].font.bold = True
        
        # 添加评审要点
        placeholder = slide.placeholders[1]
        text_frame = placeholder.text_frame
        text_frame.clear()
        
        # 分析论文中的相关文本
        relevant_sentences = analyze_innovation(paper_text, innovation_key)
        
        if relevant_sentences:
            # 限制要点数量（2-4个）
            selected_sentences = relevant_sentences[:4]
            
            for i, sentence in enumerate(selected_sentences):
                if i == 0:
                    para = text_frame.paragraphs[0]
                else:
                    para = text_frame.add_paragraph()
                
                para.text = sentence
                para.font.size = Pt(16)
                para.font.color.rgb = text_color
                para.space_after = Pt(12)
                
                # 为关键词添加高亮
                for word in innovations.values():
                    if word.split("创新")[0] in sentence:
                        para.runs[0].font.color.rgb = highlight_color
        else:
            # 如果没有找到相关内容，提供一般性描述
            para = text_frame.paragraphs[0]
            para.text = "请参考论文中关于" + innovation_title + "的具体描述"
            para.font.size = Pt(16)
            para.font.color.rgb = text_color
        
        # 添加备注
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = f"详细评审点：{innovation_title}"
    
    # 创建总结幻灯片
    summary_slide = prs.slides.add_slide(prs.slide_layouts[1])
    summary_slide.shapes.title.text = "创新点总结与评价"
    summary_slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(28)
    summary_slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = title_color
    summary_slide.shapes.title.text_frame.paragraphs[0].font.bold = True
    
    summary_placeholder = summary_slide.placeholders[1]
    summary_text_frame = summary_placeholder.text_frame
    summary_text_frame.clear()
    
    summary_points = [
        "级联双任务：创新性地将两个相关任务级联处理，可能提高模型效率与性能",
        "感知损失：引入感知损失函数，使模型学习更符合人类感知的特征表示",
        "误差预测一致性：通过误差预测一致性约束，增强模型的预测稳定性和可靠性"
    ]
    
    for i, point in enumerate(summary_points):
        if i == 0:
            para = summary_text_frame.paragraphs[0]
        else:
            para = summary_text_frame.add_paragraph()
        
        para.text = point
        para.font.size = Pt(16)
        para.font.color.rgb = text_color
        para.space_after = Pt(12)
        para.font.bold = True
    
    # 添加总体评价
    overall_para = summary_text_frame.add_paragraph()
    overall_para.text = "总体评价：本文提出的三个创新点相互补充，共同构建了一个更强大的模型框架"
    overall_para.font.size = Pt(14)
    overall_para.font.color.rgb = RGBColor(0, 100, 0)  # 深绿色
    overall_para.space_before = Pt(12)
    
    # 保存PPT
    prs.save(PPTX_PATH)
    print(f"PPT已保存到: {PPTX_PATH}")
    
    # 保存评审结果到JSON
    review_results = {
        "run_id": "vASdo9DIz1wiRpRz",
        "innovations": {
            "cascade_dual_task": {
                "title": "级联双任务",
                "description": "创新性地将两个相关任务级联处理，可能提高模型效率与性能"
            },
            "perceptual_loss": {
                "title": "感知损失",
                "description": "引入感知损失函数，使模型学习更符合人类感知的特征表示"
            },
            "error_consistency": {
                "title": "误差预测一致性",
                "description": "通过误差预测一致性约束，增强模型的预测稳定性和可靠性"
            }
        },
        "overall_evaluation": "本文提出的三个创新点相互补充，共同构建了一个更强大的模型框架"
    }
    
    json_path = OUTPUT_DIR / "review_results.json"
    with open(json_path, "w", encoding="utf-8") as f:
        json.dump(review_results, f, ensure_ascii=False, indent=2)
    
    print(f"评审结果已保存到: {json_path}")
    
    return PPTX_PATH

def main():
    # 提取论文文本
    paper_text = extract_paper_text(PDF_PATH)
    print(f"已提取论文文本，长度: {len(paper_text)} 字符")
    
    # 创建PPT
    pptx_path = create_ppt(paper_text)
    
    print(f"创新点评审报告生成完成！")
    print(f"主要文件: {pptx_path}")
    
    return pptx_path

if __name__ == "__main__":
    main()
