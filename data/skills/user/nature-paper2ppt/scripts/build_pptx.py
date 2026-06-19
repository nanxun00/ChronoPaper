#!/usr/bin/env python3
"""从 JSON 大纲生成 output/final_presentation_cn.pptx（ChronoPaper 内置）。"""
from __future__ import annotations

import argparse
import json
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt


def _add_title_slide(prs: Presentation, title: str, subtitle: str = "") -> None:
    layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title[:200]
    if subtitle and len(slide.placeholders) > 1:
        slide.placeholders[1].text = subtitle[:300]


def _add_bullet_slide(prs: Presentation, title: str, bullets: list[str]) -> None:
    layout = prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title[:200]
    body = slide.placeholders[1].text_frame if len(slide.placeholders) > 1 else None
    if body is None:
        return
    body.clear()
    for idx, bullet in enumerate(bullets[:8]):
        text = str(bullet).strip()[:500]
        if not text:
            continue
        para = body.paragraphs[0] if idx == 0 else body.add_paragraph()
        para.text = text
        para.level = 0
        para.font.size = Pt(18)


def build_from_outline(outline: dict, dest: Path) -> int:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    deck_title = str(outline.get("title") or "论文汇报").strip()
    subtitle = str(outline.get("subtitle") or "").strip()
    _add_title_slide(prs, deck_title, subtitle)

    slides = outline.get("slides") or []
    if not isinstance(slides, list):
        slides = []

    for item in slides:
        if not isinstance(item, dict):
            continue
        title = str(item.get("title") or "未命名").strip()
        bullets = item.get("bullets") or []
        if not isinstance(bullets, list):
            bullets = [str(bullets)]
        _add_bullet_slide(prs, title, [str(b) for b in bullets])

    if len(prs.slides) == 1:
        _add_bullet_slide(
            prs,
            "内容待补充",
            ["请在对话中提供更完整的论文摘要或章节要点。"],
        )

    dest.parent.mkdir(parents=True, exist_ok=True)
    prs.save(dest)
    return len(prs.slides)


def main() -> int:
    parser = argparse.ArgumentParser(description="Build PPTX from JSON outline")
    parser.add_argument("--outline", required=True, help="Path to ppt_outline.json")
    parser.add_argument(
        "--output",
        default="output/final_presentation_cn.pptx",
        help="Output PPTX path relative to skill root",
    )
    args = parser.parse_args()

    outline_path = Path(args.outline)
    if not outline_path.is_file():
        print(f"Outline not found: {outline_path}", file=sys.stderr)
        return 1

    try:
        outline = json.loads(outline_path.read_text(encoding="utf-8"))
    except (json.JSONDecodeError, OSError) as exc:
        print(f"Invalid outline: {exc}", file=sys.stderr)
        return 1

    dest = Path(args.output)
    count = build_from_outline(outline, dest)
    print(f"Wrote {dest} ({count} slides)")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
