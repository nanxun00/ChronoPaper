#!/usr/bin/env python3
"""Generate PPT from paper PDF pages for nature-paper2ppt skill."""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# --- Paths ---
SKILL_ROOT = Path(".")
RUN_ID = "y95MMvl3S15LCzGS"
OUTPUT_DIR = SKILL_ROOT / "output" / "runs" / RUN_ID
OUTPUT_DIR.mkdir(parents=True, exist_ok=True)

# PNG pages (from provided list)
PNG_DIR = SKILL_ROOT / "output" / "assets" / "figures"
page_pngs = [
    PNG_DIR / f"doi_10.48550_arxiv.2203.13474_p{p:03d}.png"
    for p in range(1, 25)
]

# --- Create presentation ---
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Helper to add a slide with a full-size image
def add_full_image_slide(png_path, title_text=None):
    slide_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(slide_layout)
    if not png_path.exists():
        # fallback: add text placeholder
        txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(11), Inches(5))
        tf = txBox.text_frame
        tf.text = f"[Image not found: {png_path.name}]"
        return slide
    # add picture scaled to slide
    slide.shapes.add_picture(str(png_path), Inches(0), Inches(0),
                             width=Inches(13.333), height=Inches(7.5))
    return slide

# --- Slide 1: Title ---
slide_layout = prs.slide_layouts[6]  # blank
slide = prs.slides.add_slide(slide_layout)
# Title background: use first page
if page_pngs[0].exists():
    slide.shapes.add_picture(str(page_pngs[0]), Inches(0), Inches(0),
                             width=Inches(13.333), height=Inches(7.5))
# Overlay text box
txBox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(3))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "PALM: Scaling Language Modeling with Pathways"
p.font.size = Pt(44)
p.font.bold = True
p.font.color.rgb = RGBColor(0x1A, 0x1A, 0x2E)
p.alignment = PP_ALIGN.CENTER
p2 = tf.add_paragraph()
p2.text = "A 540B Parameter Transformer by Google Research"
p2.font.size = Pt(24)
p2.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
p2.alignment = PP_ALIGN.CENTER

# --- Slides 2-4: Overview / Intro (pages 1-3) ---
for i in [0, 1, 2]:
    add_full_image_slide(page_pngs[i])

# --- Slides 5-7: Architecture & Training (pages 4-6) ---
for i in [3, 4, 5]:
    add_full_image_slide(page_pngs[i])

# --- Slides 8-12: Key Results (pages 7-11) ---
for i in [6, 7, 8, 9, 10]:
    add_full_image_slide(page_pngs[i])

# --- Slides 13-16: More Results & Analysis (pages 12-15) ---
for i in [11, 12, 13, 14]:
    add_full_image_slide(page_pngs[i])

# --- Slides 17-20: Downstream Tasks (pages 16-19) ---
for i in [15, 16, 17, 18]:
    add_full_image_slide(page_pngs[i])

# --- Slides 21-24: Conclusion & References (pages 20-23) ---
for i in [19, 20, 21, 22]:
    add_full_image_slide(page_pngs[i])

# --- Slide 25: Final page (page 24) ---
add_full_image_slide(page_pngs[23])

# --- Save ---
output_path = OUTPUT_DIR / "final_presentation_cn.pptx"
prs.save(str(output_path))
print(f"PPT saved to {output_path}")
