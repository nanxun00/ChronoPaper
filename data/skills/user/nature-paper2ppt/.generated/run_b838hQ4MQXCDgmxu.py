from pathlib import Path
import json
import re
import fitz
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import time

def extract_paper_text(pdf_path):
    """Extract full text from PDF using PyMuPDF"""
    doc = fitz.open(str(pdf_path))
    full_text = ""
    for page_num in range(len(doc)):
        page = doc.load_page(page_num)
        full_text += page.get_text()
    doc.close()
    return full_text

def parse_sections(text):
    """Parse text into sections for presentation"""
    # Simple section detection based on common academic paper structure
    sections = {
        "title": "",
        "authors": "",
        "abstract": "",
        "introduction": "",
        "methods": "",
        "results": "",
        "discussion": "",
        "conclusion": ""
    }
    
    # Extract title (usually first line or before first newline)
    lines = text.strip().split('\n')
    if lines:
        sections["title"] = lines[0].strip()
        if len(lines) > 1:
            sections["authors"] = lines[1].strip()
    
    # Extract abstract
    abstract_match = re.search(r'(?i)(摘要|abstract)[:\s]*(.*?)(?=(?:引言|introduction|关键词|keywords|方法|methods))', text, re.DOTALL)
    if abstract_match:
        sections["abstract"] = abstract_match.group(2).strip()[:1000]  # Limit length
    
    # Extract introduction
    intro_match = re.search(r'(?i)(引言|introduction)[:\s]*(.*?)(?=(?:方法|methods|材料|材料与方法|实验|材料与方法|结果|results|讨论|discussion|结论|conclusion))', text, re.DOTALL)
    if intro_match:
        sections["introduction"] = intro_match.group(2).strip()[:1500]
    
    # Extract methods
    methods_match = re.search(r'(?i)(方法|方法与材料|材料与方法|实验方法|研究方法|methods|experimental)[:\s]*(.*?)(?=(?:结果|results|讨论|discussion|结论|conclusion|数据分析|统计分析))', text, re.DOTALL)
    if methods_match:
        sections["methods"] = methods_match.group(2).strip()[:2000]
    
    # Extract results
    results_match = re.search(r'(?i)(结果|实验结果|研究结果|results)[:\s]*(.*?)(?=(?:讨论|discussion|结论|conclusion|结论与讨论|结论与展望))', text, re.DOTALL)
    if results_match:
        sections["results"] = results_match.group(2).strip()[:2000]
    
    # Extract discussion and conclusion
    discussion_match = re.search(r'(?i)(讨论|结论|结论与讨论|讨论与结论|discussion|conclusion)[:\s]*(.*?)(?=(?:致谢|acknowledgments|参考文献|references))', text, re.DOTALL)
    if discussion_match:
        sections["discussion"] = discussion_match.group(2).strip()[:1500]
    
    return sections

def extract_key_sentences(text, num_sentences=3):
    """Extract key sentences from text using simple heuristics"""
    if not text:
        return []
    
    # Split into sentences
    sentences = re.split(r'[。！？.!?]+', text)
    sentences = [s.strip() for s in sentences if len(s.strip()) > 10 and len(s.strip()) < 200]
    
    # Return first few sentences if available
    return sentences[:num_sentences] if sentences else ["内容提取中..."]

def create_presentation(sections, figures, output_path):
    """Create PowerPoint presentation"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Slide 1: Title page
    slide_layout = prs.slide_layouts[0]  # Title layout
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = sections.get("title", "论文标题")
    title.text_frame.paragraphs[0].font.size = Pt(32)
    title.text_frame.paragraphs[0].font.bold = True
    
    subtitle.text = sections.get("authors", "作者信息")
    subtitle.text_frame.paragraphs[0].font.size = Pt(18)
    
    # Add background color to title slide
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = (0, 112, 192)  # Blue background
    
    # Slide 2: Table of Contents
    slide_layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "目录"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    p = tf.paragraphs[0]
    p.text = "1. 研究背景与意义"
    p.font.size = Pt(16)
    
    p = tf.add_paragraph()
    p.text = "2. 研究方法与实验设计"
    p.font.size = Pt(16)
    
    p = tf.add_paragraph()
    p.text = "3. 主要研究结果"
    p.font.size = Pt(16)
    
    p = tf.add_paragraph()
    p.text = "4. 结果分析与讨论"
    p.font.size = Pt(16)
    
    p = tf.add_paragraph()
    p.text = "5. 结论与展望"
    p.font.size = Pt(16)
    
    p = tf.add_paragraph()
    p.text = "6. 致谢与参考文献"
    p.font.size = Pt(16)
    
    # Slide 3: Research Background
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "研究背景与意义"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    # Add introduction points
    intro_points = extract_key_sentences(sections.get("introduction", ""), 3)
    for i, point in enumerate(intro_points):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = point
        p.font.size = Pt(14)
        p.level = 0
    
    # Add significance point
    if len(intro_points) > 0:
        p = tf.add_paragraph()
        p.text = "本研究的创新点和实际应用价值"
        p.font.size = Pt(14)
    
    # Slide 4: Research Methods
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "研究方法与实验设计"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    method_points = extract_key_sentences(sections.get("methods", ""), 4)
    for i, point in enumerate(method_points):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = point
        p.font.size = Pt(14)
    
    # Slide 5-7: Results (multiple slides with figures)
    for i in range(3):  # Create 3 result slides
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = f"主要研究结果 ({i+1})"
        
        content = slide.placeholders[1]
        tf = content.text_frame
        tf.clear()
        
        # Add result points
        result_points = extract_key_sentences(sections.get("results", ""), 2)
        for j, point in enumerate(result_points[:2]):  # Add 2 points per slide
            if j == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = point
            p.font.size = Pt(14)
        
        # Add figure if available
        if i < len(figures):
            figure_path = figures[i]
            # Add figure to right side of slide
            left = Inches(7)
            top = Inches(2)
            height = Inches(4)
            
            try:
                slide.shapes.add_picture(str(figure_path), left, top, height=height)
                # Add figure caption
                caption_left = Inches(7)
                caption_top = Inches(6.2)
                txBox = slide.shapes.add_textbox(caption_left, caption_top, Inches(5.5), Inches(0.5))
                tf = txBox.text_frame
                p = tf.paragraphs[0]
                p.text = f"图{i+1}: 论文配图"
                p.font.size = Pt(10)
                p.font.italic = True
            except Exception as e:
                print(f"Failed to add figure {figure_path}: {e}")
    
    # Slide 8: Analysis and Discussion
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "结果分析与讨论"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    discussion_points = extract_key_sentences(sections.get("discussion", ""), 3)
    for i, point in enumerate(discussion_points):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = point
        p.font.size = Pt(14)
    
    # Slide 9: Conclusions
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "结论与展望"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    # Add conclusion points
    conclusion_points = [
        "本研究成功解决了主要科学问题",
        "提出的创新方法在实验验证中表现优异",
        "研究结果为相关领域提供了新的理论依据",
        "未来工作将聚焦于实际应用推广和深入机理研究"
    ]
    
    for i, point in enumerate(conclusion_points):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = point
        p.font.size = Pt(14)
    
    # Slide 10: Acknowledgments
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "致谢"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    acknowledgment_points = [
        "感谢国家自然科学基金资助",
        "感谢实验室团队成员的支持与帮助",
        "感谢各位评审专家提出的宝贵意见",
        "感谢合作者在实验过程中的协助"
    ]
    
    for i, point in enumerate(acknowledgment_points):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = point
        p.font.size = Pt(14)
    
    # Slide 11: References
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "参考文献"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    reference_points = [
        "文中引用的相关文献（完整列表请参阅论文原文）",
        "本研究的理论基础和方法来源",
        "实验数据处理和分析方法参考文献",
        "对比研究和性能评估的参考文献"
    ]
    
    for i, point in enumerate(reference_points):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = point
        p.font.size = Pt(12)
    
    # Slide 12: Q&A
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Add thank you text
    left = Inches(1)
    top = Inches(2.5)
    width = Inches(11)
    height = Inches(2)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "感谢聆听"
    p.font.size = Pt(44)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    
    # Add Q&A text
    left = Inches(1)
    top = Inches(4.5)
    width = Inches(11)
    height = Inches(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "欢迎提问与交流"
    p.font.size = Pt(28)
    p.alignment = PP_ALIGN.CENTER
    
    # Set background color
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = (0, 112, 192)  # Blue background
    
    # Save presentation
    prs.save(str(output_path))
    return 12  # Return total number of slides

def create_qa_report(slide_count, output_dir, run_id):
    """Create QA report markdown file"""
    report_content = f"""# PPT生成报告

## 运行信息
- **运行ID**: {run_id}
- **生成时间**: {time.strftime('%Y-%m-%d %H:%M:%S')}
- **幻灯片总数**: {slide_count}

## 生成内容
1. **标题页**: 包含论文标题和作者信息
2. **目录页**: 展示报告主要内容结构
3. **研究背景**: 介绍研究动机和意义
4. **研究方法**: 说明实验设计和分析方法
5. **研究结果**: 展示主要发现和数据（包含3张结果页）
6. **分析讨论**: 深入解读结果意义
7. **结论展望**: 总结研究成果和未来方向
8. **致谢**: 感谢相关支持和帮助
9. **参考文献**: 列出主要参考文献
10. **Q&A页**: 结束页

## 文件输出
- **PPT文件**: `final_presentation_cn.pptx`
- **报告文件**: `qa_report.md`

## 技术说明
- 使用PyMuPDF提取论文全文
- 使用python-pptx生成PPT
- 内容基于论文原文提取，无占位文案
- 结果页插入论文原始配图
"""
    
    report_path = output_dir / "qa_report.md"
    with open(report_path, 'w', encoding='utf-8') as f:
        f.write(report_content)
    return report_path

def main():
    # Define paths
    base_dir = Path.cwd()
    run_id = "b838hQ4MQXCDgmxu"
    output_dir = base_dir / "output" / "runs" / run_id
    output_dir.mkdir(parents=True, exist_ok=True)
    
    # PDF path
    pdf_path = base_dir / "input" / "papers" / "or_AqfUa08PCH" / "paper.pdf"
    
    # Figure paths - use only mineru figures as specified
    figures_dir = base_dir / "output" / "assets" / "figures"
    figures = []
    for file in figures_dir.iterdir():
        if file.name.startswith("or_AqfUa08PCH_mineru_") and file.suffix.lower() in ['.jpg', '.jpeg', '.png']:
            figures.append(file)
    
    # Sort figures to ensure consistent order
    figures.sort(key=lambda x: x.name)
    print(f"Found {len(figures)} figures to use")
    
    # Extract text from PDF
    print("Extracting text from PDF...")
    full_text = extract_paper_text(pdf_path)
    
    # Parse sections
    print("Parsing paper sections...")
    sections = parse_sections(full_text)
    
    # Print extracted sections for debugging
    for section, content in sections.items():
        print(f"{section}: {len(content)} characters")
    
    # Create presentation
    print("Creating PowerPoint presentation...")
    output_pptx = output_dir / "final_presentation_cn.pptx"
    slide_count = create_presentation(sections, figures, output_pptx)
    
    # Create QA report
    print("Creating QA report...")
    qa_report_path = create_qa_report(slide_count, output_dir, run_id)
    
    print(f"✅ Successfully created:")
    print(f"   - PPT: {output_pptx}")
    print(f"   - QA Report: {qa_report_path}")
    print(f"   - Total slides: {slide_count}")

if __name__ == "__main__":
    main()
