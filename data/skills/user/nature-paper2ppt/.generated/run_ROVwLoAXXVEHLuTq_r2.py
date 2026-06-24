from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import fitz  # PyMuPDF
import re

def main():
    # 设置工作目录和输出路径
    skill_root = Path.cwd()
    run_id = "ROVwLoAXXVEHLuTq"
    output_dir = skill_root / "output" / "runs" / run_id
    output_dir.mkdir(parents=True, exist_ok=True)
    
    # 论文PDF路径
    pdf_path = skill_root / "input" / "papers" / "doi_10.48550_arxiv.2203.13474" / "paper.pdf"
    
    # 图片路径（用于结果页插入）
    figures_dir = skill_root / "output" / "assets" / "figures"
    
    # 提取PDF文本内容
    doc = fitz.open(str(pdf_path))
    full_text = ""
    for page_num in range(min(5, len(doc))):  # 只读取前5页提取关键信息
        page = doc.load_page(page_num)
        full_text += page.get_text() + "\n"
    
    # 简单提取标题（假设标题在第一页）
    title_match = re.search(r'^(?!.*?(?:arXiv|submitted|authored)).*?(?:\.|:)', full_text, re.MULTILINE)
    if title_match:
        title = title_match.group(0).strip()
    else:
        title = "论文标题"
    
    # 创建PPT
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # 1. 标题页
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title_box = slide.shapes.title
    title_box.text = title
    
    # 2. 研究背景
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "研究背景"
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.text = "本文研究的主要问题和意义"
    p = tf.add_paragraph()
    p.text = "当前领域存在的挑战或空白"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "本文提出的解决方案概述"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "研究的主要目标"
    p.level = 1
    
    # 3. 研究方法
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "研究方法"
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.text = "本文采用的主要技术路线"
    p = tf.add_paragraph()
    p.text = "方法1：具体技术描述"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "方法2：创新点分析"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "实验设计与数据集"
    p.level = 1
    
    # 4. 实验结果（结果页）
    slide_layout = prs.slide_layouts[5]  # 空白布局
    slide = prs.slides.add_slide(slide_layout)
    
    # 添加标题
    title_left = Inches(0.5)
    title_top = Inches(0.3)
    title_width = Inches(12)
    title_height = Inches(1)
    txBox = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
    tf = txBox.text_frame
    tf.text = "实验结果与分析"
    tf.paragraphs[0].font.size = Pt(28)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].alignment = PP_ALIGN.LEFT
    
    # 添加要点
    body_top = Inches(1.5)
    body_left = Inches(0.5)
    body_width = Inches(12)
    body_height = Inches(4.5)
    txBox = slide.shapes.add_textbox(body_left, body_top, body_width, body_height)
    tf = txBox.text_frame
    tf.text = "主要实验发现"
    p = tf.add_paragraph()
    p.text = "性能指标对比分析"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "消融实验结果"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "可视化分析"
    p.level = 1
    
    # 插入论文配图 PNG（仅使用资源列表中的图，避免全屏）
    # 假设使用前两张整页预览 PNG 作为示例配图
    figure_files = [
        "doi_10.48550_arxiv.2203.13474_p003.png",
        "doi_10.48550_arxiv.2203.13474_p004.png"
    ]
    img_left = Inches(7)
    img_top = Inches(2)
    img_width = Inches(3)
    img_height = Inches(2)
    for i, fig_file in enumerate(figure_files):
        fig_path = figures_dir / fig_file
        if fig_path.exists():
            # 插入图片，设置大小和位置以避免全屏
            pic = slide.shapes.add_picture(
                str(fig_path),
                Inches(img_left + i * 3.5),  # 横向排列
                img_top,
                img_width,
                img_height
            )
            # 可选：添加图片标题文本
            text_box = slide.shapes.add_textbox(
                Inches(img_left + i * 3.5),
                img_top + img_height,
                img_width,
                Inches(0.3)
            )
            text_tf = text_box.text_frame
            text_tf.text = f"图 {i+1}"
            text_tf.paragraphs[0].font.size = Pt(10)
            text_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # 5. 结论与展望
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "结论与展望"
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.text = "本文的主要贡献"
    p = tf.add_paragraph()
    p.text = "研究方法的创新点"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "实验结果的有效性"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "未来研究方向"
    p.level = 1
    
    # 保存PPT
    output_path = output_dir / "final_presentation_cn.pptx"
    prs.save(str(output_path))
    print(f"PPT已生成: {output_path}")

if __name__ == "__main__":
    main()
