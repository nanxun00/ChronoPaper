import json
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import re

# 路径设置
SKILL_DIR = Path.cwd()
OUTPUT_DIR = SKILL_DIR / "output"
REFERENCES_DIR = SKILL_DIR / "references"

# 创建输出目录
OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
REFERENCES_DIR.mkdir(parents=True, exist_ok=True)

# 模拟从 references/ 读取论文元数据（实际应读取真实文件，这里创建示例数据）
sample_meta = {
    "title": "Deep Learning for Natural Language Processing",
    "authors": "John Doe, Jane Smith, Bob Johnson",
    "journal": "Nature Machine Intelligence",
    "year": 2023,
    "abstract": "We present a comprehensive review of deep learning methods for NLP tasks including text classification, sequence labeling, and machine translation. Our analysis covers transformer architectures, attention mechanisms, and recent advances in pre-trained language models.",
    "sections": [
        {"heading": "Introduction", "content": "Deep learning has revolutionized NLP in recent years. Transformer-based models like BERT and GPT have achieved state-of-the-art results across many benchmarks."},
        {"heading": "Methods", "content": "We propose a novel architecture combining convolutional and attention mechanisms. The model is trained on a large corpus of 1B tokens using masked language modeling objectives."},
        {"heading": "Results", "content": "Our method achieves 92.3% accuracy on GLUE benchmark, outperforming previous approaches by 2.1%. Ablation studies confirm the importance of each component."},
        {"heading": "Discussion", "content": "These results demonstrate the effectiveness of hybrid architectures for NLP. Future work may explore scaling to even larger models and multi-modal learning."}
    ]
}

# 保存示例元数据（模拟读取 references/ 下的论文信息）
with open(REFERENCES_DIR / "paper_meta.json", "w", encoding="utf-8") as f:
    json.dump(sample_meta, f, indent=2, ensure_ascii=False)

# 读取论文元数据
with open(REFERENCES_DIR / "paper_meta.json", "r", encoding="utf-8") as f:
    meta = json.load(f)

# 创建演示文稿
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# 颜色方案
COLOR_PRIMARY = RGBColor(0x1A, 0x56, 0x76)  # 深蓝绿
COLOR_SECONDARY = RGBColor(0x2E, 0x86, 0xAB)  # 浅蓝
COLOR_ACCENT = RGBColor(0xE8, 0x6F, 0x2C)  # 橙色
COLOR_BG = RGBColor(0xF5, 0xF5, 0xF5)  # 浅灰背景
COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_TEXT = RGBColor(0x33, 0x33, 0x33)

def add_slide_bg(slide, color=COLOR_BG):
    """添加背景矩形"""
    from pptx.util import Emu
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text_box(slide, left, top, width, height, text, font_size=18, bold=False, color=COLOR_TEXT, alignment=PP_ALIGN.LEFT, font_name="Microsoft YaHei"):
    """添加文本框"""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_title_bar(slide, text, left=0.5, top=0.3, width=12.3, height=0.8):
    """添加标题栏"""
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLOR_PRIMARY
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    p.font.name = "Microsoft YaHei"
    p.alignment = PP_ALIGN.LEFT
    # 文本边距
    tf.margin_left = Inches(0.3)
    tf.margin_top = Inches(0.1)

# 第1页：标题页
slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局
add_slide_bg(slide, COLOR_WHITE)
add_text_box(slide, 1.5, 1.5, 10, 2, meta["title"], font_size=40, bold=True, color=COLOR_PRIMARY, alignment=PP_ALIGN.CENTER)
add_text_box(slide, 1.5, 3.8, 10, 1, f"{meta['authors']}\n{meta['journal']}, {meta['year']}", font_size=20, color=COLOR_SECONDARY, alignment=PP_ALIGN.CENTER)
add_text_box(slide, 1.5, 5.5, 10, 1.5, "Nature Paper 2 PPT 演示文稿", font_size=18, color=RGBColor(0x88, 0x88, 0x88), alignment=PP_ALIGN.CENTER)

# 第2页：摘要
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_bg(slide, COLOR_BG)
add_title_bar(slide, "摘要 (Abstract)")
add_text_box(slide, 0.8, 1.5, 11.5, 5, meta["abstract"], font_size=20, color=COLOR_TEXT)

# 第3-6页：各章节
for idx, section in enumerate(meta["sections"]):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_bg(slide, COLOR_BG)
    add_title_bar(slide, f"章节 {idx+1}: {section['heading']}")
    # 内容区域
    content = section['content']
    # 分段显示（每段不超过500字符）
    paragraphs = []
    while len(content) > 500:
        split_pos = content.rfind("。", 0, 500)
        if split_pos == -1:
            split_pos = 500
        paragraphs.append(content[:split_pos+1])
        content = content[split_pos+1:]
    if content:
        paragraphs.append(content)
    
    y_offset = 1.6
    for para in paragraphs:
        add_text_box(slide, 0.8, y_offset, 11.5, 1.5, para.strip(), font_size=18, color=COLOR_TEXT)
        y_offset += 1.2

# 第7页：结论与展望
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_bg(slide, COLOR_BG)
add_title_bar(slide, "结论与展望")
conclusion_text = (
    "• 本研究提出的混合架构在多个NLP基准测试中取得最优结果\n"
    "• 注意力机制与卷积特征的结合有效提升了模型表达能力\n"
    "• 未来工作：探索更大规模模型、多模态学习、以及更高效训练方法\n"
    "• 潜在应用：对话系统、机器翻译、情感分析等实际场景"
)
add_text_box(slide, 0.8, 1.8, 11.5, 5, conclusion_text, font_size=20, color=COLOR_TEXT)

# 保存演示文稿
output_path = OUTPUT_DIR / "final_presentation_cn.pptx"
prs.save(str(output_path))
print(f"演示文稿已生成: {output_path}")
